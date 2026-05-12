# backend/app/services/report_generator.py
"""PDF and PowerPoint report generation for executive reports."""
from __future__ import annotations

import io
import json
import logging
from datetime import UTC, datetime
from xml.sax.saxutils import escape as xml_escape

from app.services.cosmos import CosmosRepo

logger = logging.getLogger(__name__)


class ReportGenerator:
    """Generates executive reports in PDF and PPTX formats."""

    def __init__(self, repo: CosmosRepo) -> None:
        self._repo = repo

    async def _gather_report_data(self, tenant_id: str) -> dict[str, object]:
        """Gather all data needed for executive reports."""
        summary = await self._repo.get_dashboard_summary(tenant_id)
        config = await self._repo.get_tenant_config(tenant_id)
        tenant_name = config.display_name if config else tenant_id
        return {
            "tenant_id": tenant_id,
            "tenant_name": tenant_name,
            "generated_at": datetime.now(UTC).isoformat(),
            **summary,
        }

    async def generate_executive_pdf(self, tenant_id: str) -> bytes:
        """Generate a PDF executive report.

        Uses reportlab if available, otherwise falls back to a JSON report.
        """
        data = await self._gather_report_data(tenant_id)

        try:
            from reportlab.lib.pagesizes import letter  # type: ignore[import-untyped]
            from reportlab.lib.units import inch  # type: ignore[import-untyped]
            from reportlab.platypus import (  # type: ignore[import-untyped]
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )
            from reportlab.lib.styles import getSampleStyleSheet  # type: ignore[import-untyped]
            from reportlab.lib import colors  # type: ignore[import-untyped]

            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            elements: list[object] = []

            # Title
            safe_tenant = xml_escape(str(data["tenant_name"]))
            elements.append(
                Paragraph(
                    f"Executive Security Report: {safe_tenant}",
                    styles["Title"],
                )
            )
            elements.append(Spacer(1, 0.3 * inch))
            elements.append(
                Paragraph(
                    f"Generated: {data['generated_at']}",
                    styles["Normal"],
                )
            )
            elements.append(Spacer(1, 0.5 * inch))

            # Risk summary
            elements.append(Paragraph("Risk Summary", styles["Heading2"]))
            risk_data = [
                ["Metric", "Value"],
                ["Total Identities", str(data.get("total_identities", 0))],
                [
                    "Average Risk Score",
                    f"{data.get('avg_risk_score', 0.0):.1f}",
                ],
                ["High-Risk Identities", str(data.get("high_risk_count", 0))],
                [
                    "Open Drift Alerts",
                    str(data.get("drift_alerts_open", 0)),
                ],
                [
                    "Compliance Score",
                    f"{data.get('compliance_score', 0.0):.1f}%",
                ],
            ]
            table = Table(risk_data, colWidths=[3 * inch, 2 * inch])
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2B579A")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ]
                )
            )
            elements.append(table)
            elements.append(Spacer(1, 0.5 * inch))

            # Top risky identities
            top_risky = data.get("top_risky_identities", [])
            if top_risky and isinstance(top_risky, list):
                elements.append(
                    Paragraph("Top Risky Identities", styles["Heading2"])
                )
                id_data = [["Name", "Type", "Risk Score"]]
                for identity in top_risky[:10]:
                    if isinstance(identity, dict):
                        id_data.append(
                            [
                                xml_escape(str(identity.get("display_name", ""))),
                                xml_escape(str(identity.get("identity_type", ""))),
                                f"{identity.get('risk_score', 0.0):.1f}",
                            ]
                        )
                id_table = Table(id_data, colWidths=[2.5 * inch, 1.5 * inch, 1 * inch])
                id_table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2B579A")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                            ("FONTSIZE", (0, 0), (-1, -1), 9),
                        ]
                    )
                )
                elements.append(id_table)

            doc.build(elements)
            return buffer.getvalue()

        except ImportError:
            logger.info("reportlab not installed — returning JSON report as PDF fallback")
            return json.dumps(data, indent=2, default=str).encode("utf-8")

    async def generate_executive_pptx(self, tenant_id: str) -> bytes:
        """Generate a PowerPoint executive report.

        Uses python-pptx if available, otherwise falls back to a JSON report.
        """
        data = await self._gather_report_data(tenant_id)

        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.util import Inches, Pt  # type: ignore[import-untyped]

            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = (
                f"Executive Security Report\n{data['tenant_name']}"
            )
            slide.placeholders[1].text = f"Generated: {data['generated_at']}"

            # Summary slide
            bullet_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = "Risk Summary"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = f"Total Identities: {data.get('total_identities', 0)}"
            tf.add_paragraph().text = (
                f"Average Risk Score: {data.get('avg_risk_score', 0.0):.1f}"
            )
            tf.add_paragraph().text = (
                f"High-Risk Identities: {data.get('high_risk_count', 0)}"
            )
            tf.add_paragraph().text = (
                f"Open Drift Alerts: {data.get('drift_alerts_open', 0)}"
            )
            tf.add_paragraph().text = (
                f"Compliance Score: {data.get('compliance_score', 0.0):.1f}%"
            )
            tf.add_paragraph().text = (
                f"Recommendations: {data.get('recommendations_count', 0)}"
            )

            # Top risky identities slide
            top_risky = data.get("top_risky_identities", [])
            if top_risky and isinstance(top_risky, list):
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = "Top Risky Identities"
                body = slide.placeholders[1]
                tf = body.text_frame
                for i, identity in enumerate(top_risky[:10]):
                    if isinstance(identity, dict):
                        text = (
                            f"{identity.get('display_name', 'Unknown')} "
                            f"({identity.get('identity_type', '')}) "
                            f"— Risk: {identity.get('risk_score', 0.0):.1f}"
                        )
                        if i == 0:
                            tf.text = text
                        else:
                            tf.add_paragraph().text = text

            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()

        except ImportError:
            logger.info("python-pptx not installed — returning JSON report as PPTX fallback")
            return json.dumps(data, indent=2, default=str).encode("utf-8")
